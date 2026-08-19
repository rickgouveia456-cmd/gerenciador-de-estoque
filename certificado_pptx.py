"""Gerador de Certificado de Treinamento em PPTX — Padrão Stanza."""
import io
import math
from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import OxmlElement
import lxml.etree as etree

# ── Cores ────────────────────────────────────────────────────────────────────
LARANJA  = RGBColor(0xFF, 0x6B, 0x00)
ROXO     = RGBColor(0x5C, 0x2D, 0x91)
VERDE    = RGBColor(0x00, 0x7A, 0x33)
PRETO    = RGBColor(0x00, 0x00, 0x00)
BRANCO   = RGBColor(0xFF, 0xFF, 0xFF)
CINZA    = RGBColor(0x40, 0x40, 0x40)
ROXO_CLARO = RGBColor(0x7B, 0x2F, 0xBE)
BEGE     = RGBColor(0xF5, 0xF0, 0xE8)  # fundo levemente bege

# Conteúdos programáticos por NR
CONTEUDOS = {
    'NR-18': [
        ('ASPECTOS DE SEGURANÇA DO TRABALHO:', [
            'Condições do Ambiente de Trabalho',
            'Riscos inerentes às atividades desenvolvidas',
            'Diferença entre Perigo e Risco',
            'Diferença entre Condição Insegura e Ato Inseguro',
            'Equipamentos de Proteção Coletiva',
            'Noções sobre o Programa de Gerenciamento de Riscos (PGR)',
            'Noções de Ergonomia',
        ]),
        ('NORMAS E PROCEDIMENTOS DE SEGURANÇA', [
            'Noções sobre Acidentes de Trabalho e Afastamentos',
            'Noções sobre NR 5 (CIPA)',
            'Procedimentos de Segurança em Caso de Acidentes',
        ]),
        ('FATOR PESSOAL E RELAÇÕES HUMANAS NO TRABALHO:', [
            'Respeito aos colegas, superior hierárquico e aos procedimentos internos de segurança',
            'Responsabilidade quanto aos recursos materiais disponibilizados para seu labor',
            'Função do setor de segurança do trabalho no atendimento e apoio ao trabalhador',
        ]),
    ],
    'NR-35': [
        ('CONTEÚDO PROGRAMÁTICO — TRABALHO EM ALTURA:', [
            'Normas e regulamentos aplicáveis ao trabalho em altura',
            'Análise de Risco (AR) e condições impeditivas',
            'Riscos potenciais inerentes ao trabalho em altura e medidas de prevenção e controle',
            'Sistemas, equipamentos e procedimentos de proteção coletiva',
            'EPI para trabalho em altura: seleção, inspeção, conservação e limitação de uso',
            'Acidentes típicos em trabalhos em altura',
            'Condutas em situações de emergência, incluindo noções de resgate e primeiros socorros',
        ]),
    ],
    'NR-33': [
        ('CONTEÚDO PROGRAMÁTICO — ESPAÇO CONFINADO:', [
            'Conceito de espaço confinado e riscos associados',
            'Classificação dos espaços confinados',
            'Medidas de prevenção e procedimentos de segurança',
            'Permissão de entrada e trabalho (PET)',
            'Sistemas de ventilação e monitoramento de atmosfera',
            'Equipamentos de proteção individual e coletiva',
            'Procedimentos de resgate e primeiros socorros',
        ]),
    ],
    'NR-10': [
        ('CONTEÚDO PROGRAMÁTICO — ELETRICIDADE (NR-10):', [
            'Fundamentos de eletricidade: conceitos e grandezas',
            'Riscos elétricos e medidas de controle',
            'Normas técnicas e regulamentadoras aplicáveis',
            'Proteções coletivas: aterramento, bloqueio e isolamento',
            'Equipamentos de proteção individual para eletricidade',
            'Segurança em instalações elétricas energizadas e desenergizadas',
            'Primeiros socorros em acidentes com eletricidade',
        ]),
    ],
    'NR-17': [
        ('CONTEÚDO PROGRAMÁTICO — ERGONOMIA (NR-17):', [
            'Conceito e definições — NR-17 — Ergonomia — Legislação',
            'Educação Postural no Trabalho',
            'Transporte Manual de Carga',
            'Organização no Trabalho',
            'Métodos Preventivos — Postura Adequada',
            'LER/DORT',
            'Orientação, Postura e Ginástica Laboral',
        ]),
    ],
    'NR-18 Ferramentas': [
        ('CONTEÚDO PROGRAMÁTICO — USO DE FERRAMENTAS:', [
            'Princípios de segurança na utilização de máquinas e ferramentas',
            'Operação com segurança de máquinas e ferramentas',
            'Inspeção e manutenção com segurança',
            'Sistema de bloqueio durante manutenção',
            'Manual de operação do fabricante',
            'Riscos mecânicos, elétricos e outros relevantes',
            'Método de trabalho seguro e Permissão de Trabalho',
            'Noções sobre acidentes e medidas de controle (EPC/EPI)',
        ]),
    ],
    'Brigada de Incêndio': [
        ('CONTEÚDO PROGRAMÁTICO — BRIGADA DE INCÊNDIO:', [
            'Química do fogo: triângulo e tetraedro do fogo',
            'Classificação de incêndios e agentes extintores',
            'Sistemas de detecção, alarme e combate a incêndio',
            'Operação de extintores e hidrantes',
            'Técnicas de evacuação e abandono de área',
            'Prestação de primeiros socorros',
            'Plano de emergência da obra',
        ]),
    ],
    'Primeiros Socorros': [
        ('CONTEÚDO PROGRAMÁTICO — PRIMEIROS SOCORROS:', [
            'Conceitos básicos de primeiros socorros',
            'Avaliação da cena e segurança do socorrista',
            'Hemorragias: tipos e controle',
            'Fraturas, entorses e luxações',
            'Queimaduras: classificação e atendimento',
            'Parada cardiorrespiratória e RCP',
            'Transporte de vítimas',
        ]),
    ],
}

NR_CONFIG = {
    'NR-18': {
        'nome_cert': 'Treinamento de Integração em Saúde e Segurança no Trabalho',
        'nr': 'NR 01',
        'portaria': 'Portaria nº 915 de 30/07/19',
        'carga': 4,
        'validade': 12,
    },
    'NR-35': {
        'nome_cert': 'Treinamento de Trabalho em Altura',
        'nr': 'NR 35',
        'portaria': 'Portaria nº 915 de 30/07/19',
        'carga': 8,
        'validade': 24,
    },
    'NR-33': {
        'nome_cert': 'Treinamento de Trabalho em Espaço Confinado',
        'nr': 'NR 33',
        'portaria': 'Portaria MTb 3214/78 do Ministério do Trabalho',
        'carga': 16,
        'validade': 12,
    },
    'NR-10': {
        'nome_cert': 'Treinamento de Segurança em Instalações e Serviços com Eletricidade',
        'nr': 'NR 10',
        'portaria': 'Portaria MTb 3214/78 do Ministério do Trabalho',
        'carga': 40,
        'validade': 24,
    },
    'NR-17': {
        'nome_cert': 'Noções de Ergonomia',
        'nr': 'NR 17',
        'portaria': 'Portaria 3.214/78 do Ministério do Trabalho',
        'carga': 1,
        'validade': 24,
    },
    'NR-18 Ferramentas': {
        'nome_cert': 'Treinamento para Uso de Ferramentas',
        'nr': 'NR 18',
        'portaria': 'Portaria MTb 3214/78',
        'carga': 2,
        'validade': 12,
    },
    'Brigada de Incêndio': {
        'nome_cert': 'Treinamento de Brigada de Incêndio',
        'nr': 'NR 23',
        'portaria': 'Portaria MTb 3214/78',
        'carga': 8,
        'validade': 12,
    },
    'Primeiros Socorros': {
        'nome_cert': 'Treinamento de Primeiros Socorros',
        'nr': 'NR 7',
        'portaria': 'Portaria MTb 3214/78',
        'carga': 4,
        'validade': 24,
    },
}


def _set_shape_color(shape, r, g, b, linha=False):
    """Define cor de preenchimento e remove borda."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if not linha:
        shape.line.fill.background()


def _add_rect(slide, left, top, width, height, r, g, b, rotation=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    _set_shape_color(shape, r, g, b)
    if rotation:
        shape.rotation = rotation
    return shape


def _add_line_diag(slide, left, top, width, height, r, g, b, width_pt=1.5):
    """Adiciona linha via shape de linha."""
    from pptx.util import Pt as _Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    line = slide.shapes.add_shape(9, left, top, width, height)  # rounded rect fallback
    # Usa connector ao invés
    connector = slide.shapes.add_connector(1, left, top, left+width, top+height)
    connector.line.color.rgb = RGBColor(r, g, b)
    connector.line.width = _Pt(width_pt)
    return connector


def _decoracao_canto_inferior_direito(slide, W, H):
    """Replica as setas sobrepostas roxas+laranja do canto inferior direito."""
    # Bloco laranja base (triângulo/trapézio simulado com retângulo rotacionado)
    # Retângulo laranja
    s1 = slide.shapes.add_shape(1, W - Cm(9), H - Cm(5.5), Cm(9), Cm(5.5))
    s1.fill.solid(); s1.fill.fore_color.rgb = LARANJA
    s1.line.fill.background()

    # Retângulo roxo sobreposição (seta grande)
    s2 = slide.shapes.add_shape(1, W - Cm(7.5), H - Cm(6), Cm(5), Cm(4))
    s2.fill.solid(); s2.fill.fore_color.rgb = ROXO
    s2.line.fill.background()
    s2.rotation = -15

    # Retângulo roxo claro menor (seta pequena)
    s3 = slide.shapes.add_shape(1, W - Cm(6), H - Cm(4.5), Cm(3.5), Cm(2.8))
    s3.fill.solid(); s3.fill.fore_color.rgb = ROXO_CLARO
    s3.line.fill.background()
    s3.rotation = -15

    # Retângulo laranja frente (borda laranja na seta)
    s4 = slide.shapes.add_shape(1, W - Cm(5.5), H - Cm(4), Cm(4), Cm(3.5))
    s4.fill.solid(); s4.fill.fore_color.rgb = LARANJA
    s4.line.fill.background()
    s4.rotation = -15


def _linhas_diagonais_topo_esq(slide):
    """Replica as 3 linhas diagonais roxas no canto superior esquerdo."""
    from pptx.util import Pt as _Pt
    offsets = [0, Cm(0.35), Cm(0.7)]
    for i, off in enumerate(offsets):
        try:
            conn = slide.shapes.add_connector(
                1,
                Cm(0.3) + off, Cm(0.8),
                Cm(2.2) + off, Cm(3.2)
            )
            conn.line.color.rgb = ROXO
            conn.line.width = _Pt(1.5)
        except Exception:
            pass


def _logo_seguranca_trabalho(slide):
    """Desenha o símbolo de Segurança do Trabalho com círculo verde + cruz."""
    from pptx.util import Pt as _Pt
    cx = Cm(1.8)
    cy_center = Cm(17.5)
    r = Cm(1.4)

    # Círculo externo verde escuro (borda)
    circ_ext = slide.shapes.add_shape(9, cx - r, cy_center - r, r*2, r*2)
    circ_ext.fill.solid(); circ_ext.fill.fore_color.rgb = VERDE
    circ_ext.line.color.rgb = VERDE
    circ_ext.line.width = _Pt(1)

    # Círculo interno branco
    ri = Cm(1.1)
    circ_int = slide.shapes.add_shape(9, cx - ri, cy_center - ri, ri*2, ri*2)
    circ_int.fill.solid(); circ_int.fill.fore_color.rgb = BRANCO
    circ_int.line.fill.background()

    # Cruz verde
    # Vertical
    cv = slide.shapes.add_shape(1, cx - Cm(0.18), cy_center - Cm(0.7), Cm(0.36), Cm(1.4))
    cv.fill.solid(); cv.fill.fore_color.rgb = VERDE
    cv.line.fill.background()
    # Horizontal
    ch = slide.shapes.add_shape(1, cx - Cm(0.7), cy_center - Cm(0.18), Cm(1.4), Cm(0.36))
    ch.fill.solid(); ch.fill.fore_color.rgb = VERDE
    ch.line.fill.background()

    # Texto "SEGURANÇA DO TRABALHO" em arco (simplificado como textbox curvo)
    tb = slide.shapes.add_textbox(cx - r - Cm(0.1), cy_center + Cm(0.9), r*2 + Cm(0.2), Cm(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'SEGURANÇA DO TRABALHO'
    run.font.size = Pt(4.5)
    run.font.bold = True
    run.font.color.rgb = BRANCO


def _slide_certificado(prs, participante, treinamento, empresa, cnpj, local_emissao, data_formatada):
    """Gera slide de certificado fiel ao modelo Stanza."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    W = prs.slide_width
    H = prs.slide_height

    # ── Fundo levemente bege/creme (igual ao modelo) ──────────────────────
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BEGE
    bg.line.fill.background()

    # ── Faixa laranja estreita no rodapé ──────────────────────────────────
    rodape = slide.shapes.add_shape(1, 0, H - Cm(0.3), W, Cm(0.3))
    rodape.fill.solid(); rodape.fill.fore_color.rgb = LARANJA
    rodape.line.fill.background()

    # ── Decoração canto inferior direito (setas roxo+laranja) ─────────────
    _decoracao_canto_inferior_direito(slide, W, H)

    # ── Linhas diagonais roxas canto superior esquerdo ────────────────────
    _linhas_diagonais_topo_esq(slide)

    # ── Logo Segurança do Trabalho (canto inferior esquerdo) ──────────────
    _logo_seguranca_trabalho(slide)

    # ── Título CERTIFICADO ────────────────────────────────────────────────
    tb_titulo = slide.shapes.add_textbox(Cm(3), Cm(1.5), W - Cm(6), Cm(2))
    tf = tb_titulo.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'CERTIFICADO'
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = PRETO
    run.font.name = 'Calibri'

    # ── Texto principal ───────────────────────────────────────────────────
    cfg = NR_CONFIG.get(treinamento.tipo, {})
    nome_cert = treinamento.descricao or cfg.get('nome_cert', treinamento.tipo)
    nr_ref    = treinamento.nr_referencia or cfg.get('nr', '')
    portaria  = treinamento.portaria or cfg.get('portaria', '')
    carga     = treinamento.carga_horaria or cfg.get('carga', '')
    nome_p    = (participante.colaborador or '').upper()
    cpf_p     = participante.cpf or ''
    funcao_p  = (participante.funcao or '').upper()

    tb_corpo = slide.shapes.add_textbox(Cm(2), Cm(4.0), W - Cm(11), Cm(8))
    tf = tb_corpo.text_frame
    tf.word_wrap = True

    def _run(p, text, bold=False, size=13.5):
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = PRETO
        run.font.name = 'Calibri'

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.JUSTIFY
    p.space_after = Pt(4)

    _run(p, 'Certificamos que ')
    _run(p, nome_p, bold=True)
    if cpf_p:
        _run(p, ', CPF: ')
        _run(p, cpf_p, bold=True)
    _run(p, ', na função ')
    _run(p, funcao_p, bold=True)
    _run(p, ', participou do ')
    _run(p, nome_cert, bold=True)
    _run(p, ', promovido pela empresa ')
    _run(p, empresa, bold=True)
    if cnpj:
        _run(p, ' \u2013 CNPJ: ')
        _run(p, cnpj, bold=True)
    _run(p, ', em conformidade com a ')
    if nr_ref:
        _run(p, nr_ref, bold=False)
    if portaria:
        _run(p, f', da {portaria}')
    if carga:
        _run(p, ', com carga horária de ')
        _run(p, f'{carga:02d} horas', bold=False)
    _run(p, '.')

    # ── Data ──────────────────────────────────────────────────────────────
    tb_data = slide.shapes.add_textbox(Cm(2), Cm(12.0), Cm(12), Cm(1.2))
    tf = tb_data.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f'{local_emissao}, {data_formatada}.'
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = PRETO
    run.font.name = 'Calibri'

    # ── Linha de assinatura Instrutor ─────────────────────────────────────
    inst_nome  = treinamento.responsavel or ''
    inst_cargo = treinamento.cargo_responsavel or 'Técnico em Segurança do Trabalho'
    inst_mte   = treinamento.registro_mte or ''

    # Linha horizontal instrutor
    try:
        conn = slide.shapes.add_connector(1, Cm(2.5), Cm(13.8), Cm(8.5), Cm(13.8))
        conn.line.color.rgb = PRETO
        from pptx.util import Pt as _Pt
        conn.line.width = _Pt(0.75)
    except Exception:
        pass

    tb_inst = slide.shapes.add_textbox(Cm(2), Cm(14.0), Cm(7), Cm(2.5))
    tf = tb_inst.text_frame
    tf.word_wrap = True
    linhas_inst = [
        ('Instrutor', True, 10),
        (inst_nome, False, 9),
        (inst_cargo, False, 9),
        (f'MTE: {inst_mte}' if inst_mte else '', False, 9),
    ]
    for i, (txt, bold, sz) in enumerate(linhas_inst):
        if not txt:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = PRETO
        run.font.name = 'Calibri'

    # Linha horizontal colaborador
    try:
        conn2 = slide.shapes.add_connector(1, Cm(11), Cm(13.8), Cm(17), Cm(13.8))
        conn2.line.color.rgb = PRETO
        conn2.line.width = _Pt(0.75)
    except Exception:
        pass

    tb_colab = slide.shapes.add_textbox(Cm(11), Cm(14.0), Cm(7), Cm(1.2))
    tf = tb_colab.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'Colaborador'
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = PRETO
    run.font.name = 'Calibri'

    # ── Logo stanza (texto laranja centralizado) ──────────────────────────
    tb_logo = slide.shapes.add_textbox(W//2 - Cm(4), H - Cm(3.5), Cm(8), Cm(1.8))
    tf = tb_logo.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'stanza'
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = LARANJA
    run.font.name = 'Calibri'

    return slide


def _slide_conteudo(prs, tipo_treinamento):
    """Gera slide de conteúdo programático fiel ao modelo Stanza."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    # Fundo bege
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BEGE
    bg.line.fill.background()

    # Faixa laranja rodapé
    rodape = slide.shapes.add_shape(1, 0, H - Cm(0.3), W, Cm(0.3))
    rodape.fill.solid(); rodape.fill.fore_color.rgb = LARANJA
    rodape.line.fill.background()

    # Decoração canto inferior direito
    _decoracao_canto_inferior_direito(slide, W, H)

    # Linhas diagonais
    _linhas_diagonais_topo_esq(slide)

    # Círculo Segurança do Trabalho (canto superior direito no modelo do conteúdo)
    from pptx.util import Pt as _Pt
    cx2 = W - Cm(2.5)
    cy2 = Cm(2.2)
    r2 = Cm(1.6)
    circ_e = slide.shapes.add_shape(9, cx2 - r2, cy2 - r2, r2*2, r2*2)
    circ_e.fill.solid(); circ_e.fill.fore_color.rgb = VERDE
    circ_e.line.fill.background()
    ri2 = Cm(1.2)
    circ_i = slide.shapes.add_shape(9, cx2 - ri2, cy2 - ri2, ri2*2, ri2*2)
    circ_i.fill.solid(); circ_i.fill.fore_color.rgb = BRANCO
    circ_i.line.fill.background()
    cv2 = slide.shapes.add_shape(1, cx2 - Cm(0.2), cy2 - Cm(0.8), Cm(0.4), Cm(1.6))
    cv2.fill.solid(); cv2.fill.fore_color.rgb = VERDE; cv2.line.fill.background()
    ch2 = slide.shapes.add_shape(1, cx2 - Cm(0.8), cy2 - Cm(0.2), Cm(1.6), Cm(0.4))
    ch2.fill.solid(); ch2.fill.fore_color.rgb = VERDE; ch2.line.fill.background()

    # Título CONTEÚDO PROGRAMÁTICO
    tb_tit = slide.shapes.add_textbox(Cm(2), Cm(0.8), W - Cm(7), Cm(2.2))
    tf = tb_tit.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = 'CONTEÚDO PROGRAMÁTICO'
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = PRETO
    run.font.name = 'Calibri'

    # Linha separadora vertical (barra laranja fina à esquerda do conteúdo)
    barra = slide.shapes.add_shape(1, Cm(1.8), Cm(3.5), Cm(0.12), H - Cm(5))
    barra.fill.solid(); barra.fill.fore_color.rgb = LARANJA
    barra.line.fill.background()

    # Conteúdo
    conteudo = CONTEUDOS.get(tipo_treinamento, [
        ('Conteúdo programático:', [
            'Aspectos de segurança do trabalho',
            'Riscos e medidas preventivas',
            'EPI e EPC aplicáveis',
            'Procedimentos de emergência',
        ])
    ])

    top = Cm(3.4)
    for secao_titulo, itens in conteudo:
        # Título da seção
        tb_sec = slide.shapes.add_textbox(Cm(2.2), top, W - Cm(10), Cm(0.65))
        tf = tb_sec.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = secao_titulo
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = PRETO
        run.font.name = 'Calibri'
        top += Cm(0.7)

        for item in itens:
            tb_item = slide.shapes.add_textbox(Cm(2.2), top, W - Cm(10), Cm(0.55))
            tf = tb_item.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f'- {item}'
            run.font.size = Pt(10)
            run.font.color.rgb = PRETO
            run.font.name = 'Calibri'
            top += Cm(0.52)

        top += Cm(0.35)

    # Logo stanza rodapé
    tb_logo = slide.shapes.add_textbox(W//2 - Cm(4), H - Cm(3.2), Cm(8), Cm(1.6))
    tf = tb_logo.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'stanza'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = LARANJA
    run.font.name = 'Calibri'

    return slide


def gerar_certificado_pptx(
    participantes,
    treinamento,
    empresa='STANZA INCORPORAÇÃO E CONSTRUÇÃO LTDA',
    cnpj='09.191.102/0001-06',
    local_emissao=None,
    data_formatada=None,
):
    """Gera PPTX com 2 slides por participante: Certificado + Conteúdo."""
    meses = ['Janeiro','Fevereiro','Março','Abril','Maio','Junho',
             'Julho','Agosto','Setembro','Outubro','Novembro','Dezembro']
    if not data_formatada:
        d = treinamento.data_realizacao
        data_formatada = f'{d.day:02d} de {meses[d.month-1]} de {d.year}'
    if not local_emissao:
        alm = treinamento.almoxarifado
        local_emissao = (alm.cidade or 'Local') if alm else 'Local'

    prs = Presentation()
    # Landscape 33.87 x 19.05 cm (16:9 widescreen)
    prs.slide_width  = Cm(33.87)
    prs.slide_height = Cm(19.05)

    for p in participantes:
        if not p.concluiu:
            continue
        _slide_certificado(prs, p, treinamento, empresa, cnpj,
                           local_emissao, data_formatada)
        _slide_conteudo(prs, treinamento.tipo)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
